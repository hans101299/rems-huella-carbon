from flask import Flask, request, send_file
from spire.presentation.common import *
from spire.presentation import *
from pptx import Presentation
from pptx.util import Inches
from spire.presentation import Presentation as Presentation2
import io
from flask_cors import CORS
import os

app = Flask(__name__)
CORS(app)

@app.route('/process_ppt', methods=['POST'])
def process_ppt():
    print(os.getcwd())
    data = request.json
    ppt_path = data["template"]+".pptx"
    prs = Presentation(ppt_path)

    # Supongamos que solo hay un textbox y está en la primera diapositiva
    slide = prs.slides[0]
    textbox = slide.shapes[0].text_frame

    # Realiza el reemplazo de los valores
    for key, value in data["data"].items():
        placeholder = f'{key}'
        textbox.text = textbox.text.replace(placeholder, str(value))

    # Guarda el PPT modificado
    modified_ppt_path = 'modified_ppt.pptx'
    prs.save(modified_ppt_path)

    # Convertir la primera diapositiva a una imagen
    slide_image_path = 'ToImage_0.png'
    slide_image = convert_slide_to_image(prs.slides[0])

    return send_file(slide_image_path, mimetype='image/png')

def convert_slide_to_image(slide):
    """
    Convierte una diapositiva en una imagen usando Pillow.
    """
    # Create a Presentation object
    presentation = Presentation2()

    # Load a PPT or PPTX file
    presentation.LoadFromFile("modified_ppt.pptx")

    # Loop through the slides in the presentation
    for i, slide in enumerate(presentation.Slides):

        # Specify the output file name
        fileName ="ToImage_" + str(i) + ".png"
        # Save each slide as a PNG image
        image = slide.SaveAsImage()
        image.Save(fileName)

    return image

if __name__ == '__main__':
    app.run(host='0.0.0.0',port=5002)
